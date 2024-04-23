import numpy as np
import random
import math
import torch
import matplotlib.pyplot as plt
import matplotlib.colors as mcolors
from matplotlib import cm
import scipy.io
import os
import ot
import ot.plot
import sys
import csv
import shutil
import pathlib
import argparse
import time
from statistics import mode
from pptx import Presentation
from pptx.util import Inches
from matplotlib import gridspec
sys.path.insert(0, os.path.abspath('../../'))

from wdl.bregman import barycenter
from wdl.WDL import WDL
from wdl.WDL import histRegression

from wdl.bregman import bregmanBary

from sklearn.decomposition import PCA
from sklearn.decomposition import NMF
from sklearn.manifold import TSNE
from sklearn.cluster import KMeans
from sklearn.metrics import confusion_matrix
from sklearn.neighbors import kneighbors_graph
from scipy.optimize import linear_sum_assignment

#Data file names, made as global variables for ease of use
fname = 'SalinasA_correct.mat'
matname = 'salinasA_corrected'
gtfname = 'SalinasA_gt.mat'
gtname = 'salinasA_gt' 

#Makes directory, deletes directory and subdir if it already exists
#Variables:
#path: dir name 
#override: if super important, set to false and will not replace if already exists
def dir_check(path, override=True):
    try: 
        os.mkdir(path)
    except: 
        if override: 
            shutil.rmtree(path)
            os.mkdir(path) 
        else: 
            print('Directory already exists please delete')
            exit()        
 
#Creates initial WDL directories when running experimental loop
#Variables: 
#name: directory name
def storage(name):
    path = os.path.join(os.getcwd(), name)
    dir_check(path)
    open(name + '/params.txt', 'x')

#loads in the mat file
#Variables:
#fname: file name, mat_name: reference within mat_name
def loadmat(fname, mat_name):
    root_path = os.path.dirname(os.path.dirname(os.getcwd()))
    path = os.path.join(root_path, fname)
    mat = scipy.io.loadmat(path) 
    return mat[mat_name]

#Makes all values in 2d array non-negative
#Variables:
#data: data, pad: increases lower bound for data 
def positivefy(data, pad=0):
    min = 0
    for i in range(data.shape[0]):
        for j in range(data.shape[1]):
            if data[i,j] < min:
                min = data[i,j]
    return data + np.abs(min) + pad

#Loads in data for use
#Variable:
#mode='data' or 'gt', loads in those respective files
def data_loader(mode='data', fname='SalinasA_correct.mat', matname='salinasA_corrected'):
    if mode == 'data':
        data = loadmat(fname, matname)
    elif mode == 'gt':
        data = loadmat(gtfname, gtname)
    elif mode == 'true':
        data = loadmat(fname, matname)
        return data
    data = data.reshape(data.shape[0]*data.shape[1], -1)

    if mode != 'gt': #These bands couldn't find costs for so remove them
        data = np.delete(data, [0, 32, 94], axis=1)

    return positivefy(data) #Guarantees non-negativity, a couple channels didn't hold

#This function is meant to conduct one general full experiment with outputs and everything
#Variables: 
#k: atoms 
#index: file index that has information
#train_size: size of training data
#dir_name: directory where to store everything
#reg: entropic regularizer, mu: geometric regularizer
#max_iters: num of WDL iterations, n_restarts: num of WDL restarts
#lr: learning rate 
#cost_power: In cost matrix, power used for distance
#test_size: If reconstruction, number of points
#mode: Sampling method
#n_clusters: num labels want to use, 
#label_hard: allows presetting labels used
#training_data: If non_empty, means using the training data in the file name passed in
#NOTE: mu=geometric regularizer, reg=entropic regularizer
def wdl_instance(k=2, train_size=100, dir_name='testing', reg=0.05, mu=0.1,
                 max_iters=50, n_restarts=1, lr=0.01, cost_power=1, mode='train_classes', 
                 n_clusters=2, label_hard=[], training_data='', dummy=''):
    dev = torch.device('cpu') #if torch.cuda.is_available() else torch.device("cpu")
    #torch.set_default_dtype(torch.float64) 

    if dummy == '':
        storage(dir_name) #All results saved to dir_name

    #Sets up training data, if empty will generate new random ssample
    if training_data == '':
        data = data_loader('data')
        (train_data, lst, train_classes) = sample(data, train_size, mode=mode, n_labels=n_clusters, label_hard=label_hard)
        #train_data is the data, lst is indicies in the array where data is (reshaped to 1d)
        train_index = torch.tensor(np.array(lst))
        torch.save(train_index, dir_name + '/train_index.pt')
    else:
        if dummy != '':
            train_data = torch.load(dummy + '/' + training_data)
            lst = torch.load(dummy + '/' + 'adj_index.pt')
        else:
            train_data = torch.load(training_data).float()
            train_data = train_data.to(torch.float64)
            lst = torch.load('testing_complement.pt')
    train_data = np.array(train_data)
    #Cost matrix, you can load in a file, but also Cost() generates it
    cost_mode = 'L^' + str(cost_power)
    if type(cost_power) == str:
        C = torch.load(cost_power)
    else: 
        C = Cost(cost_power)

    #Creates output file with parameters here just for tracking purposes
    with open(dir_name + '/params.txt', 'w') as file:
        file.write('cost=' + cost_mode + '\n')
        file.write('n_atoms=' + str(k) + '\n')
        file.write('mu=' + str(mu) + '\n')
        file.write('reg=' + str(reg) + '\n')
        file.write('n_restarts=' + str(n_restarts) + '\n')
        file.write('max_iter=' + str(max_iters) + '\n')
        file.write('sample size=' + str(len(lst)) + '\n')
        
    #Does WDL 
    wdl = WDL(n_atoms=k, dir=dir_name)
    train_data = train_data.T
    (weights, V_WDL) = WDL_do(dev, wdl, train_data, C, reg, mu, max_iters, lr, n_restarts)
    torch.save(V_WDL, dir_name + '/atoms.pt')
    torch.save(weights, dir_name + '/coeff.pt')

    #Visualizes learned atoms
    for i in range(0, V_WDL.shape[1]):
        plt.plot(V_WDL[:,i])
    plt.title('Learned atoms k=' + str(k) + ' mu=' + str(mu) + ' reg=' + str(reg))
    plt.savefig(dir_name + '/atoms.pdf')
    plt.clf()     

#Makes cost matrix given csv file of costs
#Variables:
#index: file index reference, cost_power: power in cost distance
def Cost(cost_power):
    vec = np.array([])
    size = 0
    file = str(os.path.dirname(os.path.dirname(os.getcwd()))) + '/salinas_costs.csv'
    with open(file) as csvfile:
        csvreader = csv.reader(csvfile, delimiter=',')
        for row in csvreader:
            vec = np.append(vec, float(row[2]))
            size += 1
    C = np.zeros((size, size))
    for i in range(0, C.shape[0]):
        for j in range(0, C.shape[1]):
            C[i, j] = abs(vec[i] - vec[j])**cost_power
    C = torch.tensor(C)
    C /= C.max()*0.1 #Divides like this to avoid numerical issues
    return C

#Does random sample given the HSI 
#Variables: 
#X: data
#Size: size of sample
#mode: Sampling modes: 
#   train_classes: Pulls from certain number of levels
#   true_random: Pulls 'size' points from anywhere
#   everything: Just gets all the data
#n_labels: how many labels want to sample from
#gt_index: File index used to pull gt labels 
#label_hard: If want to preset labels
#Data generated through call of sample under train classes
def sample(X, size, mode='train_classes', n_labels=0, label_hard=[]):
    classes = set()
    lst = set()
    gt_data = data_loader('gt')

    if mode == 'train_classes': #When want a certain number of training classes
        if len(label_hard) > 0:
            train_labels = label_hard
        else: #Will need to update labels for different images
            train_labels = random.sample([1, 10, 11, 12, 13, 14], n_labels)
        for i in range(1, len(train_labels) + 1):
            while len(lst) < i*size/len(train_labels): #Samples uniformly from each class
                val = random.randint(0, X.shape[0] - 1)
                k = gt_data[val][0]
                if k == train_labels[i-1]:
                    lst.add(val)
                    classes.add(k)
    elif 'true_random' in mode: #Only verifies data labeled, just gets random
        while len(lst) < size: 
            val = random.randint(0, gt_data.shape[0] - 1)
            if gt_data[val][0] != 0 and modif(mode, val):
                lst.add(val)
                classes.add(gt_data[val][0])
    elif 'everything' in mode: 
        for i in range(X.shape[0]):
            if gt_data[i][0] != 0 and modif(mode, i):
                lst.add(i)
                classes.add(gt_data[i][0])
    elif mode == 'restricted': 
        while len(lst) < size : 
            val = random.randint(0, gt_data.shape[0] - 1)
            if gt_data[val][0] != 0 and (val != 1045 and val != 1046) and gt_data[val] != 10:
                lst.add(val)
                classes.add(gt_data[val][0])
    train_labels = sorted(list(classes))

    if not type(lst) == list: 
        lst = list(lst)

    samp = X[lst]
    samp = samp/samp.sum(axis=1)[:,None]

    if label_hard == []:
        return (samp, lst, train_labels)
    else:
        return (samp, lst, label_hard)
    
def modif(mode, index): 
    if '++' in mode: 
        if index == 1045 or index == 1046: 
            return False
        else: 
            return True
    else:
        return True
    
#Calls wdl.fit() function in WDL class and returns atoms/weights
#More variables:
#dev: device, wdl: wdl object, init_method: WDL initialization method
#For more, on the other params, check WDL file
def WDL_do(dev, wdl, data, C, reg=0.05, mu=0.1, max_iters=100, lr=0.01, n_restarts=2, init_method='kmeans++-init'):
    weights = wdl.fit(X=torch.tensor(data).to(dev), C=C,
                init_method=init_method, loss_method="bregman",
                bary_method="bregman", reg=reg, mu=mu, max_iters=max_iters,
                max_sinkhorn_iters=50, jointOptimizer=torch.optim.Adam,
                jointOptimKWargs={"lr": lr}, verbose=True, n_restarts=n_restarts,
                log_iters=5, log=False)
    weights = weights.to("cpu")
    V_WDL = wdl.D.detach().to("cpu")
    return (weights, V_WDL)

#Normalized spectral clustering (SC w/ normalized symmetric Laplacian)
#Inputs: X: data, n_components: number of components
def spectral_cluster(X, n_components):
    Dsqrt = np.zeros((X.shape[0], X.shape[0]))
    for i in range(0, X.shape[0]):
        Dsqrt[i, i] = 1 / math.sqrt(np.sum(X[i, :]) - X[i, i])

    # inefficient symetric laplacian construction
    Lsym = np.eye(X.shape[0]) - Dsqrt @ X @ Dsqrt

    eval_sym, evec_sym = np.linalg.eigh(Lsym)
    
    V = np.zeros((X.shape[0], n_components)) #Smallest k eigenvectors
    for i in range(n_components):
        V[:,i] = evec_sym[:,i]
    #normalize row norms to 1
    V /= np.linalg.norm(V, axis=1).reshape(-1,1)

    km = KMeans(init='random', n_init='auto', n_clusters=n_components)
    km.fit(V)
    return km.labels_  

#K nearest neighbors in form of connectivity matrix
#Variables:
#W: matrix, neighbors: number of NN
#constraint: loose/tight/other, etc.
def kneighbor_weights(W, neighbors, constraint):
    W = W.T
    
    A = kneighbors_graph(W, neighbors, mode='connectivity', include_self=True)
    A = A.toarray()
    
    #What A_ij represents for each mode: 
    #None: None
    #tight: 1 if both are NN of each other, 0 otherwise
    #loose: 1 if at least one of them is NN, 0 otherwise
    if constraint == 'none': 
        return A 
    elif constraint == 'and': #exclusive/both/tight
        return np.multiply(A, A.T)
    elif constraint == 'or': #either/or
        return np.ceil((A + A.T)/2)

#Given directory name used in samples (big_sample...), gets atoms, mu, reg vals
#This might not work depending on how files are named. 
#Variable:
#path_temp: the path used
def path_convert(path_temp): 
    path_temp_k = path_temp[path_temp.find('_k=') + 1:(path_temp.find('_mu='))]
    path_temp_k = float(path_temp_k.replace('k=', ''))
    path_temp = path_temp[path_temp.find('_mu='):]
    second = path_temp.find('_') + 1
    path_temp_mu = path_temp[second: path_temp.find('_reg=')]
    path_temp_mu = float(path_temp_mu.replace('mu=', ''))
    path_temp = path_temp[second:]
    path_temp = path_temp[path_temp.find('_reg=') + 1:]
    path_temp_reg = path_temp.replace('reg=', '')
    path_temp_reg = float(path_temp_reg)
    
    return (path_temp_k, path_temp_mu, path_temp_reg)

#Clustering loop that goes through WDL results, does SC, and spatial inpainting.
#Idea is we have big parent directory and are looking through it's subdirectories. 
#Variables: 
#core_dir: Common name of sub directory where we are running these on
#NN_mode: Type of NN mentioned in kneighbors
#par_dir: Directory we are looking through to run everything
#savename: string for what you want to save the resulting NN matrix as. 
#train_mode: 'global'= using same data for everything, so loads that in. 
#recon: Will get reconstructions of training data if true. 
#For understanding, the run is clustering_loop(par_dir='/Salinas_A_experiments')
def clustering_loop(core_dir='big_sample_k=', NN_mode='or', par_dir='', 
                    savename='', train_mode='global', recon=False, savemode='HPC'):
    
    #Sets up the color map, remap gt labels for ease of use
    remap = {0: 0, 1: 1, 10: 2, 11: 3, 12: 4, 13: 5, 14: 6}
    #remap = {0: 0, 1: 1, 11: 2, 12: 3, 13: 4, 14: 5, 10: -1}
    cmap = cm.get_cmap('viridis', 7)
    new_cmap = mcolors.ListedColormap(cmap.colors)
    new_cmap.colors[0] = (1, 1, 1, 1)

    test_neighbors = [20, 25, 30, 35, 40, 45, 50] #NN we use
    
    params = np.zeros((len(test_neighbors)*1500, 5)) #Output matrix
    
    #Remaps the GT, and makes a mask for labeled points.
    (gt_data, mask) = gt_and_mask(remap)

    #If we do reconstructions, we need the cost matrix, loading it here
    if recon: 
        C = Cost(1)

    #The main loop. Iterates through neighbors, then goes through each directory
    c = 0
    for neighbors in test_neighbors: 
        for path in pathlib.Path(os.getcwd() + par_dir).iterdir():
            path_temp = str(path)
            #Checks if valid directory
            try: 
                #Gets the k, mu, and reg values. 
                #Path_convert() reliant on directory name being
                #consistent with big_sample_k=*_mu=*_reg=*
                (temp_k, temp_mu, temp_reg) = path_convert(path_temp)
            except:
                continue
            #If we are in right directory
            if core_dir in path_temp:
                #Every once in a while, NSC can fail for numerical reasons, so use this 
                try:
                    weights = torch.load(path_temp + '/coeff.pt')
                    
                    #Gets WDL reconstruction
                    if recon and neighbors == 20: 
                        atoms = torch.load(path_temp + '/atoms.pt')
                        barySolver = barycenter(C=C, method="bregman", reg=temp_reg, maxiter=100)
                        rec = np.zeros((atoms.shape[0], weights.shape[1]))
                        for i in range(weights.shape[1]):
                            rec[:,i] = barySolver(atoms, weights[:,i]).numpy().reshape(201,)
                        plt.plot(rec)
                        plt.title('WDL Reconstruction k=' + str(temp_k) + ' reg=' + str(temp_reg) + ' geom=' + str(temp_mu))
                        plt.savefig(path_temp + '/WDL_reconstruction.pdf')
                        np.save(path_temp + '/reconstructions', rec)
                        plt.close()
                    weights = weights.numpy()

                    #Cost for SC
                    for i in range(0, weights.shape[1]):
                        weights[:,i] = weights[:,i]/np.linalg.norm(weights[:,i])
                    weights = weights.T @ weights

                    W = kneighbor_weights(weights, neighbors, constraint=NN_mode)
                    labels = spectral_cluster(W, 6)
                except:
                    continue
                if train_mode != 'global':
                    index = torch.load(path_temp + '/train_index.pt').numpy()
                else:
                    index = torch.load('common_index.pt')
                linear_assignment(labels, index, remap)
                
                #Gets accuracy score
                acc = 0
                train_plot = np.zeros(83*86)
                for i in range(len(labels)):
                    t = index[i]
                    j = labels[i]
                    train_plot[t] = j
                    if gt_data[t] == j:
                        acc += 1

                #Accuracy percentage, prints out results
                acc = acc/len(labels) 
                print('atoms=' + str(temp_k), 'geom=' + str(temp_mu), 
                      'entropy=' + str(temp_reg), '| acc=' + str(acc))
                
                #Plots ground truth
                train_plot = np.reshape(train_plot, (83, 86))  
                np.save(path_temp + '/train_plot_data', train_plot)
                # if  train_mode != 'global' and neighbors == 20: 
                #     gt_grapher = np.reshape(gt_grapher, (83, 86))
                #     plt.imshow(gt_grapher, cmap=cmap)
                #     plt.savefig(path_temp + '/gt.pdf')
                #     plt.clf()
                
                #Runs spatial_NN. It can be slow, so it will only run if
                #clustering accuracy is above 60%. 
                relabel = spatial_NN(train_plot, 10, mask, run_mode='NN')
                colors = cmap(np.linspace(0, 1, cmap.N))
                new_cmap = mcolors.ListedColormap(colors)

                plt.imshow(relabel, cmap=new_cmap)
                plt.tick_params(left = False, right = False, labelleft = False, labelbottom = False, bottom = False) 
                plt.savefig(path_temp + '/Spatial_10-NN_masked_init=' + str(neighbors) + '.pdf', bbox_inches='tight')
                plt.clf()
                acc_2 = 0
                relabel = relabel.reshape(-1)
                for i in range(len(labels)):
                    t = index[i]
                    j = labels[i]
                    if gt_data[t] == relabel[t]:
                        acc_2 += 1
                acc_2 = acc_2/len(labels)
                print('Spatial relabeling accuracy | ' + str(acc_2))
                relabel = relabel.reshape((83, 86))
                p = str(round(acc_2, 2))
                plt.imshow(relabel, cmap=new_cmap)
                plt.tick_params(left = False, right = False, labelleft = False, labelbottom = False, bottom = False) 
                #plt.title('Relabeling acc=' + p)
                plt.savefig(path_temp + '/relabel_init_nn=' + str(neighbors) + '_acc=' + p + '.pdf', bbox_inches='tight')
                plt.clf()

                #Final plot
                plt.imshow(train_plot, cmap=cmap)
                plt.tick_params(left = False, right = False, labelleft = False, 
                labelbottom = False, bottom = False) 
                # plt.title('Learned labels ' + 'atoms=' + str(temp_k) + ' mu=' + str(temp_mu) 
                #             + ' reg=' + str(temp_reg) +  ' accuracy=' + str(round(acc, 2)))
                plt.savefig(path_temp + '/learned_loose=' + str(round(acc, 2)) + '_NN=' + str(neighbors) + '.pdf'
                            , bbox_inches='tight')
                plt.clf()

                #Saves the parameters and accuracy 
                params[c,:] = [temp_mu, temp_k, temp_reg, neighbors, acc]
                c += 1
    #Removes all rows that are all zeros, and then saves the matrix.
    params = params[~np.all(params == 0, axis=1)] 
    # if savemode == 'HPC':
    #     np.save(core_dir + '/NN_results_' + NN_mode, params)
    # else:
    #     np.save(os.getcwd() + par_dir + '/NN_results' + savename, params)


#Buckets: support, assuming its [0, buckets]
#a, b: uniform[a, b]
#Now it's set to be nonzero on the full support for numerical reasons
def uniform(buckets, a, b):
    X = np.zeros(buckets)
    for i in range(X.shape[0]):
        if i <= a: 
            X[i] = 1/(8*(b-a))
        if i >= b:
            X[i] = 1/(8*(b-a))
        if i > a and i < b:
            X[i] = 1/(b-a)
    X /= np.sum(X)
    return X

#Laplace(mu, b) with support [0, buckets]
def laplace(buckets, mu, b):
    Y = np.arange(0, buckets, dtype='float64')
    for i in range(Y.shape[0]):
        k = -1*abs(Y[i] - mu)
        Y[i] = math.exp(k/b)/(2*b)
    Y /= np.sum(Y, axis=0)
    return Y

#Synthetic experiments
#reg= entropic regularizer, mu=geometric regularizer, lm=want linear mixture
#To switch to other experiments
def synthetic_experiments(reg=0.001, mu=0, lm=True, dir_name='test', mode='gauss',
                          sample_size=50):
    torch.set_default_dtype(torch.float64) 
    dev = torch.device('cpu')

    dir_check(dir_name)
    samp = sample_size + 1
    size = 200 #Number of buckets 

    #Atom creation
    test = np.zeros((2, size))
    if mode == 'gauss':
        test[0,:] = ot.datasets.make_1D_gauss(size, m=50, s=5)
        test[1,:] = ot.datasets.make_1D_gauss(size, m=130, s=10)
    else: 
        test[0,:] = uniform(200, 20, 80)
        #test[1,:] = uniform(200, 130, 150)
        #test[0,:] = laplace(200, 60, 6)
        test[1,:] = laplace(200, 140, 4)

    pca_model = PCA(n_components=2)
    nmf_model = NMF(n_components=2)

    #Visualizes synthetic atoms

    if mode == 'gauss': 
        bds = [-0.025, 0.1]
    else:
        bds = [-0.025, 0.14]
    basic_pdf_plot(test.T, dir_name + '/synth_atoms_fixed.pdf', 2, ylim=bds, inv=False)

    #Creates the weights for generating the barycenters
    weights = np.zeros((samp, 2))
    for i in range(0, samp):
        k = float(1/(samp - 1))
        weights[i,:] = np.array([k*i, 1 - k*i])

    #Cost matrix for barycenters
    synth_data = np.zeros((samp, size))

    #Need to save M/M_old for when running WDL 
    M_old = cost_tensor(size)
    M = torch.tensor(np.exp(-M_old/reg)) #Kernel
    test = test.T
    
    #Linear mixture set up
    if lm: 
        test_dup = np.copy(test)
        test_dup = test_dup.T
        synth_lm = np.zeros((samp, size))
    for i in range(0, samp): #Gets synthetic data
        if lm:
            synth_lm[i,:] = weights[i,0]*test_dup[0,:] + weights[i,1]*test_dup[1,:]
        res = bregmanBary(torch.tensor(test), torch.tensor(weights[i,:]).view(-1, 1), M, reg=reg, maxiter=1000).numpy()
        synth_data[i,:] = res.reshape(res.shape[0],)

    if lm: #linear mixture visualization
        synth_lm = torch.tensor(synth_lm.T) #Plots synthetic data
        basic_pdf_plot(synth_lm, dir_name + '/linear_mixture_fixed.pdf', samp, ylim=bds, inv=True)

    #Synthetic data visualization
    basic_pdf_plot(synth_data.T, dir_name + '/synth_data_fixed.pdf', samp, ylim=bds, inv=True)
    np.save(dir_name + '/synth_data', synth_data)

    exit()

    #PCA model
    train = pca_model.fit_transform(synth_data) #PCA
    eigen = pca_model.components_
    inv = pca_model.inverse_transform(train) #PCA reconstruction

    #NMF model
    W = nmf_model.fit_transform(synth_data) #NMF 
    H = nmf_model.components_
    X = W @ H #NMF reconstruction

    #For visualizing, cycler() makes the colors line up
    #PCA components and reconstruction
    basic_pdf_plot(eigen.T, dir_name + '/PCA_evector_fixed.pdf', 2, ylim=[-0.13,0.2])
    basic_pdf_plot(inv.T, dir_name + '/PCA_reconstruct_fixed.pdf', samp, ylim=[-0.025,0.1], inv=True)

    #NMF visualization
    basic_pdf_plot(H.T, dir_name + '/NMF_components_fixed.pdf', 2, ylim=[-0.025,0.2])
    basic_pdf_plot(X.T, dir_name + '/NMF_reconstruct_fixed.pdf', samp, ylim=[-0.025,0.1], inv=True)

    #Runs WDL, as it's small, you can set n_restarts/max_iters pretty high
    wdl = WDL(n_atoms=2, dir=dir_name)
    synth_data = synth_data.T
    barySolver = barycenter(C=M_old, method="bregman", reg=reg, maxiter=1000)
    (weights, V_WDL) = WDL_do(dev, wdl, synth_data, M_old, reg, mu, n_restarts=2, max_iters=2000)
    np.save(dir_name + '/atoms', V_WDL)
    np.save(dir_name + '/weights', weights)

    #Learned atoms visualization
    #WDL initializes atoms randomly, so you might have to swap colors
    basic_pdf_plot(V_WDL, dir_name + '/learned_atoms.pdf', 2, ylim=bds)

    #Reconstruction visualization
    rec = np.zeros((V_WDL.shape[0], samp))
    for i in range(weights.shape[1]): #Gets reconstructions
        rec[:,i] = barySolver(V_WDL, weights[:,i]).numpy().reshape(size,)

    basic_pdf_plot(rec, dir_name + '/WDL_reconstruct.pdf', samp, ylim=bds)
    prs_make(dir_name)

#The code was submitted to Tufts HPC using shell scripts 
#To compartmentalize and make sure things run efficiently, rather than running
#one experiment, I set the values of mu/num atoms, then go through values of 
#reg. 
def control_loop():
    torch.set_default_dtype(torch.float64) 
    parser = argparse.ArgumentParser() #Reads in command line arguments
    parser.add_argument("--n_atoms", type=int)
    parser.add_argument("--geom", type=float)
    parser.add_argument("--track", type=int)
    parser.add_argument("--recip", type=str)
    args = parser.parse_args()
    k = args.n_atoms
    mu = args.geom
    tracker = args.track
    recip = args.recip

    if tracker == 0: 
        regs = [0.02, 0.05]
    else: 
        regs = [0.08, 0.1]

    #Shell scripts can't have non-integers, so as the <1 values of mu are
    #0.1, 0.01, and 0.001, we get around it by saying 10, 100, 1000 and taking reciprocal.
    if recip.lower() == 'true':
        mu = 1/mu 
   
    #In addition to doing WDL, this will also run the clustering loop on the results.
    for reg in regs: 
        name = 'random_sample_k=' + str(k) + '_mu=' + str(mu) + '_reg=' + str(reg)
        wdl_instance(k=k, train_size=1500, dir_name=name, reg=reg, mu=mu,
                        max_iters=400, n_restarts=2, cost_power=1, 
                        mode = 'restricted', n_clusters=6, 
                        label_hard=[1, 10, 11, 12, 13, 14], training_data='')


#Not really a way to define a Boolean one liner
def Truth(run_mode, element): 
    if run_mode == 'NN' and element == 0: 
        return True 
    elif run_mode == 'relabel' and element != 0: 
        return True 
    else:
        return False
    
#Spatial K-NN
#Now this is near exclusively run inside clustering_loop() so some of these params
#might just be 0/not used when called outside of the function.
#Inputs: 
#X: data nn: number of nn used in Spatial NN
#cmap: color map
#dir_name: directory where it's saved
#temp_k: #atoms reg: entropy mu: geom regularizer
#init_nn: #NN used in learned model 
#mask: Mask so only getting result on labeled point
def spatial_NN(X, nn, mask, run_mode='NN'): 
    #SIt's important that we only use initially labeled pixels when updating the array
    data = np.copy(X)
    c = 0

    #Outer double loop goes through every pixel, we only do stuff if labeled
    for i in range(X.shape[0]):
        for j in range(X.shape[1]):
            if Truth(run_mode, X[i,j]): 
                #3 tracking vars
                count = 0 # of nn currently added
                curr_dist = 1 #Distance from (i, j) looking at
                tracks = np.zeros(nn) #Stores results
                
                #The k-NN will be in a subgrid of the image, so we check that grid 
                #in increasing distance to get closest
                while tracks[nn-1] == 0 and curr_dist < 80: 
                    for k in range(max(i - curr_dist - 1, 0), min(i + curr_dist + 1, data.shape[0])):
                        for l in range(max(j - curr_dist - 1, 0), min(j + curr_dist + 1, data.shape[1])):
                            #So we want the vote to have a label, and check in distance range
                            if X[k,l] != 0 and euc_dist((k, l), (i, j), norm=1) == curr_dist: 
                                tracks[count] = X[k,l]
                                count += 1
                                if count == nn: #Don't continue if have the 10
                                    break
                        if count == nn: #Will double break and exit the loop
                            break
                    curr_dist += 1
                tracks = tracks[tracks != 0] #Just in case we don't get the amount
                if run_mode == 'NN': 
                    data[i, j] = mode(tracks) #Most frequent element
                elif run_mode == 'relabel':
                    if np.count_nonzero(tracks == 1) <= 1: 
                        data[i,j] = mode(tracks)
            c += 1
    data = np.multiply(data, mask) #Masks and visualizes the data
    return data

def euc_dist(X, Y, norm): #Euclidean distance for 2d
    return (np.abs(X[0] - Y[0])**norm + np.abs(X[1] - Y[1])**norm)

def grid_dist(X, Y): #Euclidean distance for 2d
    return max((np.abs(X[0] - Y[0]), np.abs(X[1] - Y[1])))

#For SalinasA, will calculate PCA/NMF for those number of atoms, just load in 
#data and pass it in to the function.
def get_pca_nmf(data):
    #All number of atoms/components used 
    arr = [2, 4, 6, 8, 10, 12, 14, 16, 18, 20, 22, 24, 26, 28, 30, 32, 34, 36]
    arr = [4]
    for k in arr:
        dir_name = 'PCA_NMF_comparisons/components=' + str(k)
        pca = PCA(n_components=k)
        nmf = NMF(n_components=k, max_iter=1000)

        train = pca.fit_transform(data) #PCA
        eigen = pca.components_
        inv = pca.inverse_transform(train) #PCA reconstruction

        W = nmf.fit_transform(data) #NMF 
        H = nmf.components_
        X = W @ H #NMF reconstruction

        plt.plot(eigen.T)
        plt.xlabel('Band number', size=15)
        plt.ylabel('Reflectance', size=15)
        plt.tick_params(axis='both', which='major', labelsize=20) # Major ticks
        plt.savefig(dir_name + '/PCA_eigenvectors_2.pdf', bbox_inches='tight')
        plt.clf()

        plt.plot(inv.T)
        plt.xlabel('Band number', size=15)
        plt.ylabel('Reflectance', size=15)
        plt.tick_params(axis='both', which='major', labelsize=20) # Major ticks
        plt.savefig(dir_name + '/PCA_reconstructions_2.pdf', bbox_inches='tight')
        plt.clf()

        plt.plot(H.T)
        plt.xlabel('Band number', size=15)
        plt.ylabel('Reflectance', size=15)
        plt.tick_params(axis='both', which='major', labelsize=20) # Major ticks
        plt.savefig(dir_name + '/NMF_components_2.pdf', bbox_inches='tight')
        plt.clf()

        plt.xlabel('Band number', size=15)
        plt.ylabel('Reflectance', size=15)
        plt.tick_params(axis='both', which='major', labelsize=20) # Major ticks
        plt.plot(X.T)
        plt.savefig(dir_name + '/NMF_reconstructions_2.pdf', bbox_inches='tight')
        plt.clf()

#Plots training data
def plot_training():
    X = torch.load('common_data.pt')
    plt.plot(X.T)
    plt.savefig('salinasA_common_data.pdf')

#Gets full spatial_NN of gt. Technically, should be roughly the same as gt. 
def gt_spatial_nn():
    cmap = cm.get_cmap('viridis', 7)
    new_cmap = mcolors.ListedColormap(cmap.colors)
    new_cmap.colors[0] = (1, 1, 1, 1)
    remap = {0: 0, 1: 1, 10: 2, 11: 3, 12: 4, 13: 5, 14: 6}

    #This remaps the GT, and makes a mask matrix. Mask is 1 if data is labeled, 
    #0 otherwise. 
    (gt_data, mask) = gt_and_mask(remap)
    spatial_NN(gt_data, 10, new_cmap, '', 0, 0, 0, 0, mask)

#Paired down version of clustering_loop() to handle comparisons for random 
#sample across data
def clustering_loop_adj(core_dir='testing_k=32', NN_mode='or', par_dir='', train_mode='local', 
                        n_labels=6):
    
    #Sets up the color map, use remap to reassign labels as matplotlib coloring
    #can be a little weird at times. 
    if n_labels == 5: 
        remap = {0: 0, 1: 1, 11: 2, 12: 3, 13: 4, 14: 5, 10: -1}
    else: 
        remap = {0: 0, 1: 1, 10: 2, 11: 3, 12: 4, 13: 5, 14: 6}
    cmap = cm.get_cmap('viridis', 7)
    new_cmap = mcolors.ListedColormap(cmap.colors)
    new_cmap.colors[0] = (1, 1, 1, 1)

    test_neighbors = [20, 25, 30, 35, 40, 45, 50] ##NN we use

    #This makes mask and remaps GT to proper labels
    (gt_data, mask) = gt_and_mask(remap)

    #The main loop. Iterates through neighbors, then goes through each directory
    #All of this is embedded in try/except to check things work
    for neighbors in test_neighbors: 
        for path in pathlib.Path(os.getcwd() + '/' + par_dir).iterdir():
            path_temp = str(path)
            #print(core_dir, path_temp)
            #If this is the right directory
            if core_dir in path_temp:
                #We use a try/except as every once in a while, normalized SC can
                #fail. 
                try:
                    weights = torch.load(path_temp + '/coeff.pt').numpy()

                    for i in range(0, weights.shape[1]):
                        weights[:,i] = weights[:,i]/np.linalg.norm(weights[:,i])
                    weights = weights.T @ weights

                    W = kneighbor_weights(weights, neighbors, constraint=NN_mode)
                    labels = spectral_cluster(W, n_labels)
                except:
                    continue
                if train_mode != 'global':
                    index = torch.load(path_temp + '/train_index.pt')
                else:
                    index = np.array(torch.load('common_index.pt'))
                labels = linear_assignment(labels, index, remap)
                #Gets accuracy score, goes through and checks if the label matches
                #the gt. 
                acc = 0
                train_plot = np.zeros(83*86)
                for i in range(len(labels)):
                    t = int(index[i])
                    j = labels[i]
                    train_plot[t] = j
                    if gt_data[t] == j:
                        acc += 1

                acc = acc/len(labels) #Accuracy
                print('Clustering NN=' + str(neighbors) + '| acc=' + str(acc))

                #Makes the ground truth plot
                train_plot = np.reshape(train_plot, (83, 86))  

                #Makes the plot of the result
                plt.imshow(train_plot, cmap=cmap)
                plt.tick_params(left = False, right = False, labelleft = False, 
                labelbottom = False, bottom = False) 
                # plt.savefig(path_temp + '/learned_Acc=' + str(round(acc, 2)) + '_NN=' + str(neighbors) + '.png'
                #             , bbox_inches='tight')
                plt.clf()

                X = spatial_NN(train_plot, 10, mask, run_mode='relabel')
                X = spatial_NN(X, 10, mask, run_mode='NN')
                colors = cmap(np.linspace(0, 1, cmap.N))
                new_cmap = mcolors.ListedColormap(colors)

                plt.imshow(X, cmap=new_cmap)
                plt.tick_params(left = False, right = False, labelleft = False, labelbottom = False, bottom = False) 
                plt.savefig(path_temp + '/Spatial_relabel_actual_init=' + str(neighbors) + '.pdf', bbox_inches='tight')
                plt.clf()
                X = np.reshape(X, (83*86))
                #Inpainting accuracy
                paint_acc = 0
                count = 0
                for i in range(83*86):
                    if X[i] != 0:
                        count += 1
                    if X[i] != 0 and X[i] == gt_data[i]:
                        paint_acc +=1 
                paint_acc = paint_acc/count
                print('IN PAINTING NN=' + str(neighbors) + '| acc=' + str(paint_acc) + '\n')

#Finds a spike and then removes and updates atoms
def spike_detect(spike_func, data, eps=0.05):
    for i in range(1, data.shape[0]):
        if data[i] > data[i - 1] + eps:
            print('spike detected: ', i)
            data = spike_func(data, i)
    return data

#Set to 0 and renormalize
def spike_correction_naive(data, idx):
    data[idx] = 0
    data /= np.sum(data)
    return data

#Redistribute mass as Gaussian
def spike_correction_naive5(data, idx):
    p = data[idx]
    data[idx] = 0

    (mu, sigma) = (np.mean(data), np.std(data))
    x = np.arange(data.shape[0], dtype=np.float64)
    gauss = np.exp(-(x-mu)**2 / (2*sigma**2))
    gauss /= gauss.sum()

    data += p*gauss
    data /= np.sum(data)
    return data

#Spike detection and removal protocol
def spike_testing(spike_func, size=5, buckets=200):
    dir_name = 'testing_samp=' + str(size)
    M = cost_tensor(buckets)
    X = np.load(dir_name + '/atoms.npy')
    weights = np.load(dir_name + '/weights.npy')
    synth = np.load(dir_name + '/synth_data.npy')

    X1 = spike_detect(spike_func, X[:,0])
    X2 = spike_detect(spike_func, X[:,1])
    X = np.array([X1, X2])

    if spike_func.__name__ == 'spike_correction_naive7': 
        wdl = WDL(n_atoms=2, dir='yo')
        wdl.fit(X=torch.tensor(synth).to('cpu'), C=M, reg=0.001, mu=0, max_iters=0,
                    jointOptimizer=torch.optim.Adam, jointOptimKWargs={"lr": 0.01})
        weights = np.zeros((2, size + 1))
        c = 0
        for i in synth:
            weights[:,c] = histRegression(X, torch.tensor(i.T), 
                                baryMethod=wdl.barycenterSolver, 
                                otMethod=wdl.OTsolver).numpy()
            c += 1
    barySolver = barycenter(C=M, method="bregman", reg=0.001, maxiter=1000)
    
    rec = np.zeros((200, size + 1))
    for i in range(weights.shape[1]): #Gets reconstructions
        rec[:,i] = barySolver(torch.tensor(X.T), torch.tensor(weights[:,i])).numpy().reshape(buckets,)
    basic_pdf_plot(rec, dir_name + '/updated_recon_naive.pdf', 6, ylim=[-0.025, 0.1])
    np.save(dir_name + '/up_recons_naive.npy', rec)

    basic_pdf_plot(X.T, dir_name + '/updated_atoms_naive.pdf', 2, ylim=[-0.025, 0.15])
    np.save(dir_name + '/up_atoms_relearn.npy', rec)

def synth_exp():
    sizes = [5, 10, 20, 25, 50, 100]
    for size in sizes:
        dir_name = 'uniform_laplace/uni_laplace_testing=' + str(size)
        synthetic_experiments(reg=0.001, lm=False, dir_name=dir_name, mode='laplace', 
                            sample_size=size)
        
#Basic plot saved as pdf under viridis
def basic_pdf_plot(data, savename, v_size, ylim=[0,1], inv=False):
    if inv :
        plt.gca().set_prop_cycle(plt.cycler('color', plt.cm.viridis(np.linspace(0, 1, v_size))))
    else:
        plt.gca().set_prop_cycle(plt.cycler('color', plt.cm.viridis(np.linspace(1, 0, v_size))))
    #plt.tick_params(axis='both', which='major', labelsize=22) # Major ticks
    plt.plot(data)
    plt.ylim(ylim[0], ylim[1])
    #plt.savefig(savename, bbox_inches='tight')
    plt.savefig(savename)
    plt.clf()

def tsne(dir_name, comp=2, n_labels=6):
    if n_labels == 6: 
        remap = {0: 0, 1: 1, 10: 2, 11: 3, 12: 4, 13: 5, 14: 6}
    elif n_labels == 5: 
        remap = {0: 0, 1: 1, 11: 2, 12: 3, 13: 4, 14: 5, 10: -1}

    cmap = cm.get_cmap('viridis', 7)
    new_cmap = mcolors.ListedColormap(cmap.colors)
    new_cmap.colors[0] = (1, 1, 1, 1)
    X = torch.load(dir_name + '/coeff.pt')
    X = X.T
    embed = TSNE(n_components=comp, learning_rate='auto', init='random', 
                 perplexity=25).fit_transform(X)
    index = torch.load(dir_name + '/train_index.pt')

    km = KMeans(init='k-means++', n_init='auto', n_clusters=n_labels)
    gt_data = data_loader('gt')
    for i in range(gt_data.shape[0]):
        gt_data[i] = remap.get(gt_data[i][0])
    
    gt_temp = np.zeros(len(index))
    for i in range(index.shape[0]):
        gt_temp[i] = int(gt_data[index[i]])

    km.fit(embed)
    labels = km.labels_    
    labels = linear_assignment(labels, index, remap)

    #Gets accuracy score
    acc = 0
    for i in range(len(labels)):
        t = index[i]
        j = labels[i]
        if gt_data[t] == j:
            acc += 1
    #Accuracy percentage, prints out results
    acc = acc/len(labels) 
    print(dir_name[93:] + ' components accuracy=' + str(acc))
    plt.scatter(embed[:,0], embed[:,1], c=gt_temp, cmap='viridis')
    plt.title('TSNE initial data')
    plt.savefig(dir_name + '/tsne_true.pdf', bbox_inches='tight')
    plt.clf()
    plt.scatter(embed[:,0], embed[:,1], c=labels, cmap='viridis')
    plt.title('TSNE labeling post kmeans acc=' + str(round(acc, 2)))
    plt.savefig(dir_name + '/tsne_kmeans')
    plt.clf()

def linear_assignment(learned_labels, index, remap):
    label_set = sorted(list(set(learned_labels)))

    gt_data = data_loader('gt')
    for i in range(gt_data.shape[0]):
        gt_data[i] = remap.get(gt_data[i][0])
    
    gt_temp = np.zeros(len(index))
    for i in range(index.shape[0]):
        gt_temp[i] = int(gt_data[index[i]])

    #Need to remap the resulting SC labels to the correct ones
    gt_labs = np.array(list(set(gt_temp)))
    for i in range(len(learned_labels)):
        learned_labels[i] = gt_labs[label_set.index(learned_labels[i])]
    #Linear assignment to match clusters
    confusion = confusion_matrix(gt_temp, learned_labels)
    cost_final = -confusion + np.max(confusion)
    (res1, res2) = linear_sum_assignment(cost_final)

    #Remaps for visualization 
    temp2 = list(gt_labs[res2])
    for i in range(0, len(learned_labels)):
        learned_labels[i] = temp2.index(learned_labels[i]) + 1
    return learned_labels

#This is just fancy kmeans with extra steps
def tsne_means(dir_name, tol=0.05):
    X = torch.load(dir_name + '/coeff.pt')
    Y = torch.load(dir_name + '/train_index.pt')
    gt = data_loader('gt')

    X = X.T
    remap = {0: 0, 1: 1, 10: 2, 11: 3, 12: 4, 13: 5, 14: 6}
    cmap = cm.get_cmap('viridis', 7)
    new_cmap = mcolors.ListedColormap(cmap.colors)
    new_cmap.colors[0] = (1, 1, 1, 1)
    for i in range(gt.shape[0]):
        gt[i] = remap.get(gt[i][0])

    gt_train = gt[Y] #Might do nothing
    c = np.zeros(gt_train.shape[0])

    base = [np.zeros(5) for i in range(0, 6)]
    for i in range(0, 6):
        j = 0
        while j < 5:
            val = random.randint(0, 7137)
            if (gt[val][0] - 1) == i and not (val in base[i]):
                base[i][j] = val
                c[np.where(Y == gt[val][0])[0][0]] = gt[val][0]
                j += 1
    embed = TSNE(n_components=2, learning_rate='auto', init='random', 
                perplexity=25).fit_transform(X)
    plt.title('Tsne embedding labels match')
    plt.scatter(embed[:,0], embed[:,1], c=gt_train, cmap='viridis')
    plt.savefig(dir_name + '/tsne_embed')
    plt.clf()

    #This is just kmeans with extra steps
    pts_labeled = 0
    while pts_labeled < 7137: 
        for i in range(0, embed.shape[0]):
            exit = False
            for temp in base:
                if Y[i] in temp:
                    exit = True
                if exit:
                    break
            if exit:
                continue
            else: 
                mean = np.zeros(6)
                for k in range(0, len(base)): 
                    dist = 0
                    for l in range(0, base[k].shape[0]):
                        y = (np.where(Y == int(base[k][l])))[0][0]
                        dist += np.linalg.norm(embed[int(y)] - embed[i])
                    mean[k] = dist/base[k].shape[0]
                mean /= mean.sum()
                if np.min(mean) < tol:
                    new_l = np.argmin(mean) 
                    c[i] = new_l
                    base[int(new_l)] = np.append(base[int(new_l)], Y[i])
                    pts_labeled += 1
                    if pts_labeled == 7137:
                        break
        plt.scatter(embed[:,0], embed[:,1], c=c)
        plt.show()
        plt.clf()
        tol += 0.025

def prs_make(dir_name):
    iter_track = [50*i for i in range(0, 50)]
    prs = Presentation()
    for i in range(0, 50):
        try: 
            os.remove(dir_name + '/iter=' + str(iter_track[i]) + '.pdf')
        except: 
            x = 2
        path = dir_name + '/iter=' + str(iter_track[i]) + '.png'
        left = Inches(2)
        top = Inches(1)
        width = Inches(7)
        height = Inches(6)
        try: 
            prs.slides.add_slide(prs.slide_layouts[0])
            pic = slide.shapes.add_picture(path, left, top, width, height)
        except: 
            break 
    prs.save(dir_name + '/time_lapse_pres.pptx')

#This function given index builds complement of it 
def set_complement(dir_name):
    torch.set_default_dtype(torch.float64) 

    X = list(torch.load(dir_name + '/train_index.pt').numpy())
    GT = data_loader('gt')

    X_C = list()
    for i in range(GT.shape[0]):
        if not (i in X) and (GT[i][0] != 0) and (i != 1045 and i != 1046):
            X_C.append(i)

    return torch.tensor(np.array(X_C))

def gt_and_mask(remap, mask_marker=0): 
    gt_data = data_loader('gt')
    mask = np.zeros(83*86)
    for i in range(gt_data.shape[0]):
        gt_data[i] = remap.get(gt_data[i][0])
        if gt_data[i] != 0 or gt_data[i] != mask_marker:
            mask[i] = 1
    mask = np.reshape(mask, (83, 86))
    return (gt_data, mask)

def cost_tensor(size):
    M = ot.utils.dist0(size)
    M /= M.max()
    return torch.tensor(M)

def ssl_reg(point_find, rho=0.01, init_points=1): 
    random.seed(10)
    start = time.time()
    dir_name = 'random_sample_tests/random_samp_32/random_sample_k=32_mu=0.0001_reg=0.08'
    X = torch.load(dir_name + '/coeff.pt').numpy()
    idx = torch.load(dir_name + '/train_index.pt').numpy()
    num_points = idx.shape[0]

    #Needed initialization information
    remap = {0: 0, 1: 1, 10: 2, 11: 3, 12: 4, 13: 5, 14: 6}
    labels = [1, 2, 3, 4, 5, 6]
    all_points = set()
    point_grouping = [set() for i in range(0, 6)]

    (gt_data, mask) = gt_and_mask(remap)

    #Gets initial points in each class
    for i in range(1, len(labels) + 1):
        while len(list(point_grouping[i - 1])) < init_points: 
            val = random.randint(0, gt_data.shape[0] - 1)
            if gt_data[val][0] == labels[i-1] and (val in idx):
                point_grouping[i - 1].add(val)
                all_points.add(val)
    dup_point = all_points.copy()
    #Gets initial associated weights
    point_grouping = [list(point_grouping[i]) for i in range(0, 6)]
    weights = [np.zeros((32,init_points)) for i in range(0, 6)]
    for i in range(0, 6): 
        weights[i][:,0] = X[:,np.where(idx == point_grouping[i][0])[0][0]]
    
    weight_d_tracker = [np.zeros((num_points, init_points)) for i in range(0, 6)]
    #On scale, spatial-norm dominates, so meant to scale things down
    
    s_norm = 100
    
    #SCORE UPDATE LOOP
    #UPDATE 2, add only certain points
    acc = 0
    error = 0.2
    tot_lab = 0
    while (error < 0.21) and (tot_lab != (X.shape[1] - len(dup_point))):
    #for marking in range(0, point_find): 
        m_score = np.zeros((7138, 6))
        for i in range(0, num_points):
            z = idx[i]
            if not (z in all_points):
                #Grid location of image
                k1 = math.floor(z/86) 
                k2 = z - k1*86

                #This calculates the score vector
                score = np.zeros(6)
                for j in range(0, 6): 
                    (d1, d2) = (0, 0)
                    size = weights[j].shape[0]
                    for k in range(0, weights[j].shape[1]): 
                        if weight_d_tracker[j][i, k] == 0: 
                            dist = np.linalg.norm(X[:,i] - weights[j][:,k])
                            weight_d_tracker[j][i, k] = dist
                            d1 += dist
                        else:
                            d1 += weight_d_tracker[j][i, k]
                        test_loc = point_grouping[j][k]
                        x_d = math.floor(test_loc/86)
                        y_d = test_loc - x_d*86
                        d2 += (np.abs(k1-x_d) + np.abs(k2-y_d))
                    
                    d2 /= s_norm
                    score[j] = (d1 + rho*d2)/size
                score /= np.sum(score)
                m_score[z,:] = score

        #Finds the minimum point
        non_zero = np.any(m_score != 0, axis=1)
        m_score = m_score[non_zero]
        idx_2 = np.zeros(7138)
        for i in range(0, idx.shape[0]):
            idx_2[idx[i]] = int(idx[i])
        idx_2 = idx_2[non_zero]

        vals = np.argmin(m_score, axis=1)
        p_lab = 0
        for i in range(0, vals.shape[0]):
            score = vals[i] + 1
            #if m_score[i, vals[i]] < error:
            point_grouping[score - 1].append(int(idx_2[i]))
                # all_points.add(int(idx_2[i]))
                # weights[vals[i]] = np.append(weights[vals[i]], X[:,i].reshape(-1, 1), axis=1)
                # weight_d_tracker[vals[i]] = np.append(weight_d_tracker[vals[i]], np.zeros((num_points, 1)), axis=1)
                # p_lab += 1

        print('Points labeed: ', p_lab, error)
        error += 0.01
        tot_lab += p_lab

        # coloring = np.zeros(7138)
        # for i in range(0, 7138): 
        #     for j in range(0, len(point_grouping)): 
        #         if i in point_grouping[j]: 
        #             coloring[i] = j + 1
        #     if i in dup_point:
        #         coloring[i] = 7

        # coloring = np.reshape(coloring, (83, 86))
        # cmap_temp = cm.get_cmap('viridis', 8)
        # cmap_temp2 = cm.get_cmap('viridis', 7)
        # new_cmap = mcolors.ListedColormap(cmap_temp.colors) 
        # fix = mcolors.ListedColormap(cmap_temp2.colors) 
        
        # for i in range(0, cmap_temp2.colors.shape[0]):
        #     new_cmap.colors[i] = fix.colors[i]
        # new_cmap.colors[0] = (1, 1, 1, 1)
        # new_cmap.colors[7] = (1, 0, 0, 1)

        # plt.imshow(coloring, cmap=new_cmap)
        # plt.title('Accuracy = ' + str(round(acc, 2)))
        # plt.show()
        #Pixel by pixel variant
        # row_index, col_index = np.unravel_index(np.argmin(m_score), m_score.shape)
        # m_row = m_score[row_index,:]
        # m_lab_idx = np.argmin(m_row)
        # label = m_lab_idx + 1
        # idx_2 = int(idx_2[row_index])

        # #NOW THIS COMPLETES ONE ITERATION: 
        # point_grouping[m_lab_idx].append(idx_2)
        # weights[m_lab_idx] = np.append(weights[m_lab_idx], X[:,row_index].reshape(-1, 1), axis=1)
        # weight_d_tracker[m_lab_idx] = np.append(weight_d_tracker[m_lab_idx], np.zeros((num_points, 1)), axis=1)

        # all_points.add(idx_2)
        # if label == gt_data[idx_2][0]:
        #     acc +=1
        # print('Prediction:', label, 'True label:', gt_data[idx_2][0], marking)
            
    for i in range(0, len(point_grouping)): 
        for j in point_grouping[i]: 
            if (not (j in dup_point)) and (i + 1) == gt_data[j][0]: 
                acc += 1
    acc = acc/X.shape[1]
    print(rho, init_points,' TOTAL ACCUARCY: ',  acc)
    coloring = np.zeros(7138)
    for i in range(0, 7138): 
        for j in range(0, len(point_grouping)): 
            if i in point_grouping[j]: 
                coloring[i] = j + 1
        if i in dup_point:
            coloring[i] = 7

    coloring = np.reshape(coloring, (83, 86))
    cmap_temp = cm.get_cmap('viridis', 8)
    cmap_temp2 = cm.get_cmap('viridis', 7)
    new_cmap = mcolors.ListedColormap(cmap_temp.colors) 
    fix = mcolors.ListedColormap(cmap_temp2.colors) 
    
    for i in range(0, cmap_temp2.colors.shape[0]):
        new_cmap.colors[i] = fix.colors[i]
    new_cmap.colors[0] = (1, 1, 1, 1)
    new_cmap.colors[7] = (1, 0, 0, 1)

    plt.imshow(coloring, cmap=new_cmap)
    plt.title('Accuracy = ' + str(round(acc, 2)))
    plt.show()
    #plt.savefig('ssl_test2/rho=' + str(rho) + '_init_points=' + str(init_points) + '.png')
    plt.clf()
    return (acc, list(dup_point))
    

def virid_modify():
    cmap = cm.get_cmap('viridis', 7)
    new_cmap = mcolors.ListedColormap(cmap.colors)
    new_cmap.colors[0] = (1, 1, 1, 1)
    return new_cmap

def cooking(dir_name):
    N = 100
    X = torch.load(dir_name + '/coeff.pt').numpy()
    index = torch.load(dir_name + '/train_index.pt').numpy()

    Y = np.zeros((X.shape[0] + 2, X.shape[1]))
    for i in range(X.shape[0]):
        Y[i,:] = X[i,:]

    for i in range(X.shape[1]): 
        k1 = math.floor(index[i]/86)
        k2 = index[i] - k1*86
        Y[X.shape[0], i] = k1/N
        Y[X.shape[0] + 1, i] = k2/N
    
    remap = {0: 0, 1: 1, 10: 2, 11: 3, 12: 4, 13: 5, 14: 6}
    gt_data = data_loader('gt')
    for i in range(gt_data.shape[0]):
        gt_data[i] = remap.get(gt_data[i][0])
    
    gt_temp = np.zeros(len(index))
    for i in range(index.shape[0]):
        gt_temp[i] = int(gt_data[index[i]])

    Y = Y.T
    # embed = TSNE(n_components=2, learning_rate='auto', init='random', 
    #              perplexity=25).fit_transform(Y)
    # plt.scatter(embed[:,0], embed[:,1], c=gt_temp, cmap='viridis')
    # plt.savefig('embed_cooking')
    # plt.clf()

    weights = Y @ Y.T
    for neighbors in [20, 25, 30, 35, 40, 45, 50]: 
        W = kneighbor_weights(weights, neighbors, constraint='or')
        labels = spectral_cluster(W, 6)
        labels = linear_assignment(labels, index, remap)
        #Gets accuracy score, goes through and checks if the label matches
        #the gt. 
        acc = 0
        train_plot = np.zeros(83*86)
        for i in range(len(labels)):
            t = int(index[i])
            j = labels[i]
            train_plot[t] = j
            if gt_data[t] == j:
                acc += 1
        acc = acc/len(labels) #Accuracy
        print('N: ', N, 'Neighbors: ', neighbors, 'Acc:', acc)
        train_plot = np.reshape(train_plot, (83, 86))  

        #Makes the plot of the result
        cmap = virid_modify()
        plt.imshow(train_plot, cmap=cmap)
        plt.show()
        plt.clf()


if __name__ == "__main__":
    print('Imports complete') #Just some default stuff, change dev if using gpus
    # torch.set_default_dtype(torch.float64) 
    # np.set_printoptions(suppress=True) #Indices are 1045, 1046 of outliers


    for path in pathlib.Path(os.getcwd() + '/Salinas_A_experiments').iterdir():
        path = str(path)
        if 'big_sample_k=32_' in path: 
            X = torch.load(path + '/atoms.pt')
            plt.plot(X)
            plt.savefig(path + '/atoms_fixed.pdf', bbox_inches='tight')
            plt.clf()
